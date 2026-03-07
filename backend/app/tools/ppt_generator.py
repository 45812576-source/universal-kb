"""PPT generation builtin tool.

Input params:
{
  "title": "演示标题",
  "slides": [
    {"title": "幻灯片1标题", "content": "正文内容", "layout": "title_content"},
    ...
  ]
}

Output: {"file_id": "xxx.pptx", "download_url": "/api/files/xxx.pptx"}
"""
from __future__ import annotations

import os
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

_UPLOAD_DIR = os.environ.get("UPLOAD_DIR", "./uploads")
_GENERATED_DIR = Path(_UPLOAD_DIR) / "generated"


def execute(params: dict) -> dict:
    """Generate a PPTX file and return file_id + download_url."""
    _GENERATED_DIR.mkdir(parents=True, exist_ok=True)

    title_text = params.get("title", "演示文稿")
    slides_data = params.get("slides", [])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title_text
    if slide.placeholders[1]:
        slide.placeholders[1].text = ""

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_data in slides_data:
        slide = prs.slides.add_slide(content_layout)
        slide_title = slide_data.get("title", "")
        slide_content = slide_data.get("content", "")

        if slide.shapes.title:
            slide.shapes.title.text = slide_title
        if len(slide.placeholders) > 1:
            tf = slide.placeholders[1].text_frame
            tf.word_wrap = True
            tf.text = slide_content

    file_id = f"{uuid.uuid4().hex}.pptx"
    file_path = _GENERATED_DIR / file_id
    prs.save(str(file_path))

    return {
        "file_id": file_id,
        "download_url": f"/api/files/{file_id}",
        "filename": f"{title_text}.pptx",
        "message": f"PPT已生成，共{len(slides_data) + 1}张幻灯片",
    }

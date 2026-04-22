"""Extract plain text from uploaded documents."""
from __future__ import annotations

import base64
import html
import os
import re
import tempfile
import zipfile
from dataclasses import dataclass
from xml.etree import ElementTree as ET


def convert_pdf_to_docx(pdf_path: str, timeout: int = 120) -> str:
    """将 PDF 转换为 DOCX，返回生成的 DOCX 文件路径。

    使用 pdf2docx 库保留排版、表格、图片。
    调用者负责清理返回的临时文件。
    timeout 秒后自动终止，防止设计类 PDF 卡死。
    """
    import subprocess
    import sys
    import json as _json

    docx_path = pdf_path.rsplit(".", 1)[0] + ".docx"

    # 在子进程中运行 pdf2docx，防止卡死阻塞主进程
    script = (
        "import sys, json; "
        "from pdf2docx import Converter; "
        "cv = Converter(sys.argv[1]); "
        "cv.convert(sys.argv[2]); "
        "cv.close(); "
        "print(json.dumps({'ok': True}))"
    )
    try:
        result = subprocess.run(
            [sys.executable, "-c", script, pdf_path, docx_path],
            capture_output=True, text=True, timeout=timeout,
        )
        if result.returncode != 0:
            stderr = (result.stderr or "").strip()[-500:]
            raise RuntimeError(f"pdf2docx 转换失败 (exit {result.returncode}): {stderr}")
    except subprocess.TimeoutExpired:
        raise RuntimeError(f"pdf2docx 转换超时 ({timeout}s)，该 PDF 可能包含复杂图形排版")

    if not os.path.exists(docx_path):
        raise RuntimeError("pdf2docx 未生成 DOCX 文件")

    return docx_path


def _call_kimi_vision(image_path: str) -> str:
    """Call Kimi vision API (via 百炼 Coding Plan) to OCR/describe an image."""
    import openai

    api_key = os.environ.get("BAILIAN_API_KEY", "")
    if not api_key:
        raise ValueError("BAILIAN_API_KEY 环境变量未设置")

    with open(image_path, "rb") as f:
        b64_data = base64.b64encode(f.read()).decode("utf-8")

    ext = os.path.splitext(image_path)[1].lower().lstrip(".")
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "webp": "image/webp", "bmp": "image/bmp", "gif": "image/gif"}
    mime = mime_map.get(ext, "image/png")

    client = openai.OpenAI(
        api_key=api_key,
        base_url="https://coding.dashscope.aliyuncs.com/apps/openai/v1",
    )
    resp = client.chat.completions.create(
        model="kimi-k2.5",
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64_data}"}},
                    {"type": "text", "text": "请详细描述这张图片的内容，包括所有文字、数据、图表信息。"},
                ],
            }
        ],
        max_tokens=2048,
    )
    return resp.choices[0].message.content or ""


def _call_ark_vision(image_path: str) -> str:
    """Call ARK-compatible vision endpoint when configured."""
    import openai

    api_key = os.environ.get("ARK_API_KEY", "")
    base_url = os.environ.get("ARK_BASE_URL", "").strip()
    model = os.environ.get("ARK_VISION_MODEL", "").strip()
    if not api_key or not base_url or not model:
        raise ValueError("ARK vision 环境变量未完整配置")

    with open(image_path, "rb") as f:
        b64_data = base64.b64encode(f.read()).decode("utf-8")

    ext = os.path.splitext(image_path)[1].lower().lstrip(".")
    mime_map = {"jpg": "image/jpeg", "jpeg": "image/jpeg", "png": "image/png",
                "webp": "image/webp", "bmp": "image/bmp", "gif": "image/gif"}
    mime = mime_map.get(ext, "image/png")

    client = openai.OpenAI(api_key=api_key, base_url=base_url)
    resp = client.chat.completions.create(
        model=model,
        messages=[
            {
                "role": "user",
                "content": [
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{b64_data}"}},
                    {"type": "text", "text": "请做 OCR 与版面理解，尽量逐段输出图片中的原文，不要编造。"},
                ],
            }
        ],
        max_tokens=4096,
    )
    return resp.choices[0].message.content or ""


def _call_vision_ocr(image_path: str) -> str:
    provider = os.environ.get("KNOWLEDGE_VISION_PROVIDER", "").strip().lower()
    errors: list[str] = []
    providers = [provider] if provider else ["bailian", "ark"]

    for candidate in providers:
        try:
            if candidate == "ark":
                return _call_ark_vision(image_path)
            return _call_kimi_vision(image_path)
        except Exception as exc:
            errors.append(f"{candidate}:{exc}")
    raise ValueError("视觉 OCR 调用失败: " + " | ".join(errors))


@dataclass
class ExtractionResult:
    text: str
    mode: str
    error: str | None = None


_XLSX_NS = {
    "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "rel": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
_XLSX_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"


def _xlsx_col_to_index(cell_ref: str) -> int:
    letters = "".join(ch for ch in cell_ref if ch.isalpha())
    if not letters:
        return 0
    index = 0
    for ch in letters.upper():
        index = index * 26 + (ord(ch) - 64)
    return max(index - 1, 0)


def _xlsx_shared_strings(zf: zipfile.ZipFile) -> list[str]:
    if "xl/sharedStrings.xml" not in zf.namelist():
        return []
    root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
    values: list[str] = []
    for item in root.findall("main:si", _XLSX_NS):
        texts = [node.text or "" for node in item.findall('.//main:t', _XLSX_NS)]
        values.append("".join(texts))
    return values


def _xlsx_cell_value(cell: ET.Element, shared_strings: list[str]) -> str:
    cell_type = cell.attrib.get("t")
    if cell_type == "inlineStr":
        texts = [node.text or "" for node in cell.findall('.//main:t', _XLSX_NS)]
        return "".join(texts)

    value_node = cell.find("main:v", _XLSX_NS)
    if value_node is None or value_node.text is None:
        return ""

    raw_value = value_node.text
    if cell_type == "s":
        try:
            shared_idx = int(raw_value)
        except (TypeError, ValueError):
            return raw_value
        if 0 <= shared_idx < len(shared_strings):
            return shared_strings[shared_idx]
        return raw_value
    if cell_type == "b":
        return "TRUE" if raw_value == "1" else "FALSE"
    return raw_value


def _extract_xlsx_rows_fallback(file_path: str) -> list[tuple[str, list[list[str]]]]:
    import posixpath

    with zipfile.ZipFile(file_path) as zf:
        shared_strings = _xlsx_shared_strings(zf)
        workbook_root = ET.fromstring(zf.read("xl/workbook.xml"))
        rels_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
        rel_map = {
            rel.attrib.get("Id"): rel.attrib.get("Target", "")
            for rel in rels_root.findall(f'{{{_XLSX_REL_NS}}}Relationship')
        }

        sheets: list[tuple[str, list[list[str]]]] = []
        for sheet in workbook_root.findall("main:sheets/main:sheet", _XLSX_NS):
            sheet_name = sheet.attrib.get("name") or "Sheet"
            rel_id = sheet.attrib.get(f'{{{_XLSX_NS["rel"]}}}id')
            target = rel_map.get(rel_id)
            if not target:
                continue
            sheet_path = posixpath.normpath(posixpath.join("xl", target))
            sheet_root = ET.fromstring(zf.read(sheet_path))

            rows: list[list[str]] = []
            for row in sheet_root.findall("main:sheetData/main:row", _XLSX_NS):
                values: list[str] = []
                last_index = -1
                for cell in row.findall("main:c", _XLSX_NS):
                    ref = cell.attrib.get("r", "")
                    cell_index = _xlsx_col_to_index(ref) if ref else last_index + 1
                    while len(values) < cell_index:
                        values.append("")
                    values.append(_xlsx_cell_value(cell, shared_strings))
                    last_index = cell_index
                if any((value or "").strip() for value in values):
                    rows.append(values)
            sheets.append((sheet_name, rows))

        return sheets


def _extract_xlsx_text_fallback(file_path: str) -> str:
    parts: list[str] = []
    for sheet_name, rows in _extract_xlsx_rows_fallback(file_path):
        rendered_rows: list[str] = []
        for row in rows:
            if any((cell or "").strip() for cell in row):
                rendered_rows.append("	".join(row))
        if rendered_rows:
            parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rendered_rows))
    return "\n\n".join(parts)


def _looks_like_meaningful_pdf_text(text: str) -> bool:
    compact = "".join(text.split())
    return len(compact) >= 40


def _pdf_pages_to_images(file_path: str) -> list[str]:
    import fitz

    temp_dir = tempfile.mkdtemp(prefix="pdf-ocr-")
    image_paths: list[str] = []
    doc = fitz.open(file_path)
    try:
        for idx, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            out = os.path.join(temp_dir, f"page-{idx}.png")
            pix.save(out)
            image_paths.append(out)
    finally:
        doc.close()
    return image_paths


def _extract_pdf_with_vision(file_path: str) -> ExtractionResult:
    image_paths = _pdf_pages_to_images(file_path)
    if not image_paths:
        return ExtractionResult(text="", mode="pdf_fallback", error="pdf_page_render_failed")

    texts: list[str] = []
    errors: list[str] = []
    try:
        for image_path in image_paths:
            try:
                text = _call_vision_ocr(image_path).strip()
                if text:
                    texts.append(text)
            except Exception as exc:
                errors.append(str(exc))
    finally:
        for image_path in image_paths:
            try:
                os.unlink(image_path)
            except OSError:
                pass
        try:
            os.rmdir(os.path.dirname(image_paths[0]))
        except OSError:
            pass

    combined = "\n\n".join(texts).strip()
    if combined:
        return ExtractionResult(text=combined, mode="pdf_vision_ocr", error=None)
    return ExtractionResult(
        text="",
        mode="pdf_fallback",
        error=("vision_ocr_failed: " + " | ".join(errors))[:500] if errors else "vision_ocr_failed",
    )


def _strip_html_tags(raw_html: str) -> str:
    if not raw_html:
        return ""
    cleaned = re.sub(r"(?is)<(script|style).*?>.*?</\1>", " ", raw_html)
    cleaned = re.sub(r"(?i)<br\s*/?>", "\n", cleaned)
    cleaned = re.sub(r"(?i)</p\s*>", "\n", cleaned)
    cleaned = re.sub(r"(?i)</div\s*>", "\n", cleaned)
    cleaned = re.sub(r"(?s)<[^>]+>", " ", cleaned)
    cleaned = html.unescape(cleaned)
    cleaned = re.sub(r"\r\n?", "\n", cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned.strip()


def extract_text_result(file_path: str) -> ExtractionResult:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        import pdfplumber

        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        joined = "\n\n".join(texts).strip()
        if _looks_like_meaningful_pdf_text(joined):
            return ExtractionResult(text=joined, mode="pdf_text")
        return _extract_pdf_with_vision(file_path)

    return ExtractionResult(text=extract_text(file_path), mode="text")


def _transcribe_funasr(audio_path: str) -> str:
    """Transcribe audio using local FunASR model."""
    import re
    from funasr import AutoModel

    model = AutoModel(
        model="iic/SenseVoiceSmall",
        vad_model="fsmn-vad",
        punc_model="ct-punc",
        disable_update=True,
    )
    res = model.generate(input=audio_path, batch_size_s=300)
    if not res:
        return ""
    text = res[0].get("text", "")
    # Remove SenseVoice emotion/language tags like <|zh|><|NEUTRAL|>
    text = re.sub(r"<\|[^|]+\|>", "", text).strip()
    return text


_FOE_MAP_PROMPT = """\
你是严谨的信息分析师。对下面的文本片段做摘要，并对每个关键论断用FOE三测试标注：
- [F] 事实：可通过公认方法验证真假
- [O] 观点：两个同等知情者可能合理分歧，需注明发布方和隐含假设
- [E] 证据：为支撑某判断而引用，需注明质量层级（一级=审计/官方 二级=权威二手 三级=专家报告 四级=类比先例 五级=观点用作证据）

输出格式：
## Chunk摘要
### 核心事实
- [F] ...
### 关键观点
- [O] ... —— 发布方：... | 隐含假设：...
### 证据链
- [E] ... —— 质量层级：X级 | 支撑观点：...
### 本段主旨
[一句话，明确区分F/O成分]

文本片段：
{text}"""

_FOE_REDUCE_PROMPT = """\
下面是一篇长文各段落的FOE标注摘要。请合并为最终结构化摘要。

合并规则：
1. 事实去重，相同事实保留更精确版本
2. 观点聚合，标注支撑强度（强≥2条一/二级证据；中=1条高质量或多条三级；弱=仅四/五级；无=纯断言）
3. 冲突的事实/观点明确标注矛盾，保留双方
4. 将分散证据重新关联到对应观点

输出格式：
## 一、事实摘要
（按主题聚类，去重）

## 二、观点图谱
| 观点 | 发布方 | 支撑强度 | 隐含假设 |
|------|--------|----------|---------|

## 三、证据评估
- 整体论证质量：
- 证据缺口：
- 观点伪装为证据的环节：

## 四、信息缺口
（应讨论但未讨论的方面）

## 五、全文摘要
（3-5句，每句标注[F]或[O]）

---
各段落摘要：
{text}"""


def foe_summarize(raw_text: str, llm_cfg: dict) -> str:
    """对超长文本执行 LangChain MapReduce + FOE 三测试摘要。

    返回 FOE 结构化摘要字符串（用于聊天上下文）。
    原始文本由调用方单独保存到知识库，数据不丢失。

    llm_cfg: llm_gateway.resolve_config(db, slot_key) 返回的 dict，
             包含 api_base / api_key / model_id 等字段。
    """
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_openai import ChatOpenAI
    from langchain_core.messages import HumanMessage

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=4000,
        chunk_overlap=500,
        separators=["\n\n", "\n", "。", "，", " "],
    )
    chunks = splitter.split_text(raw_text)

    llm = ChatOpenAI(
        base_url=llm_cfg.get("api_base", "https://api.deepseek.com/v1"),
        api_key=llm_cfg.get("api_key", ""),
        model=llm_cfg.get("model_id", "deepseek-chat"),
        temperature=0.1,
        max_tokens=2000,
    )

    # Map: 每个 chunk 独立做 FOE 摘要
    chunk_summaries = []
    for chunk in chunks:
        prompt = _FOE_MAP_PROMPT.format(text=chunk)
        resp = llm.invoke([HumanMessage(content=prompt)])
        chunk_summaries.append(resp.content)

    # Reduce: 合并所有 chunk 摘要
    combined = "\n\n---\n\n".join(chunk_summaries)
    reduce_prompt = _FOE_REDUCE_PROMPT.format(text=combined)
    final = llm.invoke([HumanMessage(content=reduce_prompt)])
    return final.content


def extract_html(file_path: str) -> str:
    """从文件提取 HTML 格式内容，保留富文本格式（标题/加粗/列表/表格等）。
    用于前端云文档编辑器渲染。不支持的格式回退到 extract_text() 包装。"""
    ext = os.path.splitext(file_path)[1].lower()

    if ext in (".docx",):
        import mammoth
        _style_map = "\n".join([
            "p[style-name='Heading 1'] => h1:fresh",
            "p[style-name='Heading 2'] => h2:fresh",
            "p[style-name='Heading 3'] => h3:fresh",
            "p[style-name='标题 1'] => h1:fresh",
            "p[style-name='标题 2'] => h2:fresh",
            "p[style-name='标题 3'] => h3:fresh",
        ])
        with open(file_path, "rb") as f:
            result = mammoth.convert_to_html(f, style_map=_style_map)
            return result.value

    elif ext in (".md",):
        import markdown as md_lib
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return md_lib.markdown(text, extensions=["tables", "fenced_code", "nl2br", "sane_lists"])

    elif ext in (".html", ".htm"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    elif ext in (".pptx",):
        # PPTX: wrap each slide as a section with headings
        from pptx import Presentation
        prs = Presentation(file_path)
        html_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text.strip()
                    if t:
                        slide_texts.append(t)
            if slide_texts:
                title = slide_texts[0]
                body = slide_texts[1:] if len(slide_texts) > 1 else []
                html_parts.append(f"<h2>{title}</h2>")
                for line in body:
                    html_parts.append(f"<p>{line}</p>")
        return "\n".join(html_parts)

    elif ext in (".xlsx", ".xls", ".csv"):
        # Tabular data → HTML table
        plain = extract_text(file_path)
        lines = plain.strip().split("\n")
        html_parts = []
        in_table = False
        for line in lines:
            if line.startswith("[Sheet:"):
                if in_table:
                    html_parts.append("</tbody></table>")
                    in_table = False
                html_parts.append(f"<h3>{line.strip('[]')}</h3>")
            elif "\t" in line:
                if not in_table:
                    html_parts.append("<table><tbody>")
                    in_table = True
                cells = line.split("\t")
                html_parts.append("<tr>" + "".join(f"<td>{c}</td>" for c in cells) + "</tr>")
            elif line.strip():
                if in_table:
                    html_parts.append("</tbody></table>")
                    in_table = False
                html_parts.append(f"<p>{line}</p>")
        if in_table:
            html_parts.append("</tbody></table>")
        return "\n".join(html_parts)

    elif ext in (".txt",):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            text = f.read()
        return "\n".join(f"<p>{line or '<br>'}</p>" for line in text.split("\n"))

    elif ext in (".pdf",):
        result = extract_text_result(file_path)
        if not result.text:
            return ""
        return "\n".join(
            f"<p>{html.escape(line) if line else '<br>'}</p>"
            for line in result.text.split("\n")
        )

    else:
        # 其他格式（PDF/图片/音频等）：纯文本包装为 <p>
        try:
            plain = extract_text(file_path)
            return "\n".join(f"<p>{line or '<br>'}</p>" for line in plain.split("\n"))
        except ValueError:
            return ""


def extract_text(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    elif ext == ".pdf":
        import pdfplumber
        texts = []
        with pdfplumber.open(file_path) as pdf:
            for page in pdf.pages:
                t = page.extract_text()
                if t:
                    texts.append(t)
        return "\n\n".join(texts)

    elif ext in (".docx",):
        from docx import Document
        from docx.oxml.ns import qn

        doc = Document(file_path)
        parts = []

        def iter_block_items(parent):
            """按文档顺序迭代段落和表格。"""
            from docx.oxml.ns import qn as _qn
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            parent_elm = parent.element.body if hasattr(parent, 'element') else parent
            for child in parent_elm.iterchildren():
                if child.tag == qn('w:p'):
                    yield Paragraph(child, parent)
                elif child.tag == qn('w:tbl'):
                    yield Table(child, parent)

        def extract_table(table):
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                # 合并相邻重复单元格（跨列合并会重复出现）
                deduped = []
                for c in cells:
                    if not deduped or c != deduped[-1]:
                        deduped.append(c)
                row_text = "\t".join(c for c in deduped if c)
                if row_text:
                    rows.append(row_text)
            return "\n".join(rows)

        for block in iter_block_items(doc):
            from docx.table import Table
            from docx.text.paragraph import Paragraph
            if isinstance(block, Paragraph):
                t = block.text.strip()
                if t:
                    parts.append(t)
            elif isinstance(block, Table):
                t = extract_table(block)
                if t:
                    parts.append(t)

        return "\n\n".join(parts)

    elif ext in (".pptx",):
        from pptx import Presentation
        prs = Presentation(file_path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text.strip())
        return "\n\n".join(t for t in texts if t)

    elif ext in (".md",):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()

    elif ext in (".html", ".htm"):
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            return _strip_html_tags(f.read())

    elif ext in (".xlsx", ".xls"):
        if ext == ".xlsx":
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, data_only=True)
                parts = []
                for sheet in wb.worksheets:
                    rows = []
                    for row in sheet.iter_rows(values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            rows.append("	".join(cells))
                    if rows:
                        parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
                return "\n\n".join(parts)
            except Exception:
                return _extract_xlsx_text_fallback(file_path)

        import openpyxl
        wb = openpyxl.load_workbook(file_path, data_only=True)
        parts = []
        for sheet in wb.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    rows.append("	".join(cells))
            if rows:
                parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
        return "\n\n".join(parts)

    elif ext in (".csv",):
        import csv
        rows = []
        with open(file_path, "r", encoding="utf-8", errors="replace") as f:
            reader = csv.reader(f)
            for row in reader:
                rows.append("\t".join(row))
        return "\n".join(rows)

    elif ext in (".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"):
        return _call_kimi_vision(file_path)

    elif ext in (".mp3", ".wav", ".m4a", ".ogg", ".flac"):
        return _transcribe_funasr(file_path)

    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            "Supported: .txt .pdf .docx .pptx .md .html .htm .xlsx .xls .csv "
            ".jpg .jpeg .png .webp .bmp .gif .mp3 .wav .m4a .ogg .flac"
        )
